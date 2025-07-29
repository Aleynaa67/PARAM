from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()

# Başlık slaytı
slide_0 = prs.slides.add_slide(prs.slide_layouts[0])
slide_0.shapes.title.text = "PARAM Ödeme ve İade Entegrasyon Projesi"
slide_0.placeholders[1].text = "Aleyna Akılıç"

# Yardımcı fonksiyon: madde listesi slaytı
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullet_points[0]
    for point in bullet_points[1:]:
        content.text += f"\n{point}"
    return slide

# Slaytlar
add_bullet_slide("🎯 Proje Amacı", [
    "ParamPOS altyapısıyla ödeme, iade ve dekont işlemlerini gerçekleştirmek.",
    "SOAP web servisleri kullanarak API entegrasyonu sağlamak.",
    "Estetik ve kullanıcı dostu bir arayüz geliştirmek."
])

add_bullet_slide("⚙️ Kullanılan Teknolojiler", [
    "Python (Flask framework)",
    "HTML, CSS, JavaScript",
    "XML & SOAP Web Servisleri",
    "ReportLab & SMTP ile PDF ve e-posta gönderimi"
])

add_bullet_slide("💳 Ödeme İşlemi", [
    "Kredi kartı bilgileri formdan alınır.",
    "İşlem hash’i oluşturulur (SHA1 + Base64).",
    "SOAP ile ödeme API’sine XML gönderilir.",
    "Başarılı işlem sonrası detaylar kullanıcıya gösterilir."
])

add_bullet_slide("↩️ İade / İptal İşlemi", [
    "Sipariş ID ve tutar ile işlem alınır.",
    "SOAP üzerinden API’ye 'IADE' veya 'IPTAL' gönderilir.",
    "Sonuç mesajı kullanıcıya aktarılır."
])

add_bullet_slide("📨 Dekont Gönderimi", [
    "İşlem ID ve e-posta alınır.",
    "PDF dekont oluşturulur (ReportLab).",
    "Kullanıcının e-posta adresine otomatik gönderim yapılır."
])

add_bullet_slide("✨ Ek Özellikler", [
    "#ParamİleGüvende etiketi (sağ alt köşe)",
    "Yıldızlı animasyonlu arka plan",
    "Mobil uyumlu HTML / CSS yapısı",
    "Kullanıcı dostu başarı / hata mesajları"
])

add_bullet_slide("🔐 Güvenlik", [
    "Kart verisi saklanmaz, yalnızca işlem için kullanılır.",
    "SHA1 ve GUID ile güvenli işlem akışı.",
    "SOAP servisi ile kimlik doğrulama yapılır."
])

add_bullet_slide("📈 Kazanımlarım", [
    "SOAP/XML ile API kullanımını öğrendim.",
    "Gerçek ödeme servisleriyle çalışma deneyimi kazandım.",
    "Frontend ve backend entegrasyonunu uyguladım.",
    "Kullanıcı deneyimi odaklı arayüzler tasarladım."
])

add_bullet_slide("🧪 Demo (İsteğe Bağlı)", [
    "Ödeme işlemi gerçekleştirme",
    "İade veya iptal örneği",
    "Dekont oluşturup e-posta ile gönderme"
])

# Sunumu kaydet
prs.save("Param_Proje_Sunumu.pptx")
